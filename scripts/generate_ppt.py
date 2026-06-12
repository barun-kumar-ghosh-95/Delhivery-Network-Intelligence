import collections 
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Define a helper function to set slide title and format
    def set_slide_title(slide, text):
        title_shape = slide.shapes.title
        title_shape.text = text
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(36)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x2d, 0x37, 0x48)
                
    def add_bullet_points(slide, points):
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # clear default bullet
        for i, point in enumerate(points):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = point
            p.space_after = Pt(14)
            for run in p.runs:
                run.font.size = Pt(24)

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Optimizing Delivery ETAs with Graph-Based Network Intelligence"
    subtitle.text = "Analytics Competition Submission\nDelhivery Network Optimization"

    # Slide 2: Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "1. The Problem")
    add_bullet_points(slide, [
        "Inaccurate ETA predictions cause SLA breaches and customer dissatisfaction.",
        "Baseline OSRM (Open Source Routing Machine) systematically underestimates actual travel time by nearly 2x.",
        "Facility processing delays and bottlenecks are ignored by traditional models.",
        "Goal: Build an intelligent, network-aware ETA prediction model."
    ])

    # Slide 3: Dataset
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "2. The Dataset & Engineering")
    add_bullet_points(slide, [
        "140,000+ logistics trips across the Delhivery network.",
        "Feature Engineering:",
        "  - Temporal (hour of day, day of week, rush hour flags)",
        "  - Spatial (haversine distance, city pairs)",
        "  - Aggregations (segment-level rollups for actual time)",
        "Data Quality: Addressed extreme outliers (500+ hour delays) using IQR thresholds."
    ])

    # Slide 4: Graph Construction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "3. Graph Construction")
    add_bullet_points(slide, [
        "Modeled the entire logistics network as a Directed Graph using NetworkX.",
        "Nodes: 1,657 Facilities (Hubs & Centers).",
        "Edges: 2,783 Corridors (Trip routes).",
        "Edge Weights: Average historical delay on the corridor.",
        "Extracted topological features for each facility: PageRank, Betweenness Centrality, Degree Centrality."
    ])

    # Slide 5: Graph Features (Network Vis & Graph Features)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "4. Why Graph Topology Matters")
    add_bullet_points(slide, [
        "Node2Vec Embeddings: 32-dimensional vectors capturing the structural neighborhood of each facility.",
        "Community Detection (Louvain): Grouped facilities into 61 regional sub-networks.",
        "Impact: Tabular models treat facilities as isolated categorical variables. Graph features expose the 'ripple effect' of congestion.",
    ])

    # Slide 6: Bottleneck Identification
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "5. Identifying Bottlenecks")
    add_bullet_points(slide, [
        "Computed a composite 'Impact Score' combining:",
        "  - PageRank (Network Importance)",
        "  - Throughput (Volume)",
        "  - SLA Breach % (Historical Delay Probability)",
        "Result: Identified Top 5 critical nodes responsible for >15% of all network delays.",
        "Actionable Insight: Focus capacity expansion at these exact nodes."
    ])

    # Slide 7: Baseline vs Graph Models
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "6. Model Comparison")
    add_bullet_points(slide, [
        "Baseline XGBoost (Tabular Only): MAE 72.31 mins",
        "Node2Vec + XGBoost: MAE 64.00 mins",
        "Graph + LightGBM (Final): MAE 62.09 mins | R² 0.91",
        "Error Reduction: 14.1% improvement over tabular baseline.",
        "Accuracy @ 15% Error Margin: Reached 67.4%."
    ])

    # Slide 8: FTL vs Carting Advisor
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "7. FTL vs Carting Intelligence")
    add_bullet_points(slide, [
        "Different vehicle types exhibit different delay patterns.",
        "FTL (Full Truck Load) vs. Carting (Intracity vans).",
        "Built an interactive advisor to recommend the optimal vehicle type per corridor.",
        "Corridors with >2.0x delay ratio are flagged for immediate vehicle re-allocation."
    ])

    # Slide 9: Business Impact & Command Center
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "8. Business Impact")
    add_bullet_points(slide, [
        "Developed 'Network Operations Command Center' Dashboard.",
        "Real-time monitoring of Corridors, Bottlenecks, and ETAs.",
        "Estimated Revenue at Risk tracking based on SLA breaches and throughput.",
        "Projected SLA improvement: 20-30%.",
        "Estimated Revenue Recovery: Rs. 2-5 Cr annually."
    ])

    # Slide 10: Recommendations & Future Scope
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "9. Recommendations & Future Scope")
    add_bullet_points(slide, [
        "90-Day Plan: Deploy ETA correction model on top 50 corridors.",
        "12-Month Plan: Capacity expansion at Top 5 Bottlenecks.",
        "Future Scope: Implement Graph Neural Networks (GraphSAGE/GAT) for inductive learning.",
        "Future Scope: Dynamic, real-time edge weight updates based on live traffic."
    ])

    prs.save('presentation.pptx')
    print("Successfully generated presentation.pptx")

if __name__ == '__main__':
    create_presentation()
